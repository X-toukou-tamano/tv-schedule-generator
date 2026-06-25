import os

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor

from datetime import datetime
from zoneinfo import ZoneInfo
# ==================================================
# 確定値
# ==================================================

# X座標
NIGHT_X = Cm(1.128)
DAY_X = Cm(17.637)

# サイズ
BOX_WIDTH = Cm(15.239)

BOX_HEIGHT_SINGLE = Cm(5.165)
BOX_HEIGHT_MULTI = Cm(3.379)

# Y座標
Y_POSITIONS = {
    1: [Cm(7.595)],
    2: [
        Cm(5.906),
        Cm(11.071)
    ],
    3: [
        Cm(4.828),
        Cm(8.680),
        Cm(12.532)
    ]
}

# 背景色
NIGHT_FILL = RGBColor(51, 86, 147)
DAY_FILL = RGBColor(255, 255, 255)

# 文字色
NIGHT_FONT = RGBColor(255, 255, 255)
DAY_FONT = RGBColor(0, 0, 0)


# ==================================================
# 表示文字作成
# ==================================================

def build_display_text(
    name,
    grade,
    status
):
    """
    場名 + グレード + 日数
    """

    name = str(name).strip()
    grade = str(grade).strip()
    status = str(status).strip()

    return (
        f"{name}"
        f"　     "
        f"{grade}"
        f"      　"
        f"{status}"
    )


# ==================================================
# 四角作成
# ==================================================

def add_schedule_box(
    slide,
    x,
    y,
    width,
    height,
    text,
    fill_color,
    font_color
):

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x,
        y,
        width,
        height
    )

    # 背景色
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    # 枠線なし
    shape.line.fill.background()

    # テキスト
    tf = shape.text_frame
    tf.clear()

    tf.vertical_anchor = (
        MSO_VERTICAL_ANCHOR.MIDDLE
    )

    p = tf.paragraphs[0]

    run = p.add_run()
    run.text = text

    font = run.font
    font.name = "Rockwell"
    font.size = Pt(28)
    font.bold = True
    font.color.rgb = font_color

    return shape


# ==================================================
# メイン
# ==================================================

def create_powerpoint(
    day_events,
    night_events
):

    template_path = os.path.join(
        os.getcwd(),
        "src",
        "元データ.pptx"
    )

    if not os.path.exists(template_path):
        raise FileNotFoundError(
            f"テンプレートが存在しません: {template_path}"
        )

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # 最大3場

    if len(day_events) > 3:
        raise Exception(
            f"デイが3場を超えています: {len(day_events)}"
        )

    if len(night_events) > 3:
        raise Exception(
            f"ナイターが3場を超えています: {len(night_events)}"
        )

    # ==================================================
    # ナイター
    # ==================================================

    night_count = len(night_events)

    if night_count > 0:

        height = (
            BOX_HEIGHT_SINGLE
            if night_count == 1
            else BOX_HEIGHT_MULTI
        )

        for idx, event in enumerate(
            night_events
        ):

            text = build_display_text(
                event.get("name", ""),
                event.get("grade", ""),
                event.get("status", "")
            )

            add_schedule_box(
                slide=slide,
                x=NIGHT_X,
                y=Y_POSITIONS[night_count][idx],
                width=BOX_WIDTH,
                height=height,
                text=text,
                fill_color=NIGHT_FILL,
                font_color=NIGHT_FONT
            )

    # ==================================================
    # デイ
    # ==================================================

    day_count = len(day_events)

    if day_count > 0:

        height = (
            BOX_HEIGHT_SINGLE
            if day_count == 1
            else BOX_HEIGHT_MULTI
        )

        for idx, event in enumerate(
            day_events
        ):

            text = build_display_text(
                event.get("name", ""),
                event.get("grade", ""),
                event.get("status", "")
            )

            add_schedule_box(
                slide=slide,
                x=DAY_X,
                y=Y_POSITIONS[day_count][idx],
                width=BOX_WIDTH,
                height=height,
                text=text,
                fill_color=DAY_FILL,
                font_color=DAY_FONT
            )

    # ==================================================
    # 保存
    # ==================================================

    upload_dir = os.path.join(
        os.getcwd(),
        "uploads"
    )

    os.makedirs(
        upload_dir,
        exist_ok=True
    )

    today_str = (
        datetime.now(
            ZoneInfo("Asia/Tokyo")
        )
        .strftime(
            "%m%d"
        )
    )

    output_path = os.path.join(
        upload_dir,
        f"場内放映予定{today_str}.pptx"
    )

    prs.save(output_path)

    print(
        f"[SUCCESS] {output_path}"
    )

    return output_path
